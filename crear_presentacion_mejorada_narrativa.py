#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para crear presentación con narrativa mejorada y formato corregido.
Enfoque: Corredores Mineros y Turísticos Complementarios

MEJORAS:
1. Narrativa clara: Problema → Solución → Método → Resultados
2. Formato de actividades corregido (sin saltos de línea)
3. Colores con mejor contraste
4. Jerarquía visual mejorada
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def crear_presentacion():
    """Crear presentación profesional con narrativa mejorada"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colores con mejor contraste
    COLOR_AZUL_PRINCIPAL = RGBColor(0, 74, 151)      # Azul GIZ
    COLOR_VERDE_EXITO = RGBColor(0, 135, 104)        # Verde GIZ
    COLOR_NARANJA_ENFASIS = RGBColor(255, 127, 39)   # Naranja destacado
    COLOR_MORADO_ACCION = RGBColor(111, 66, 193)     # Morado para actividades
    COLOR_ROJO_PROBLEMA = RGBColor(220, 53, 69)      # Rojo para problemas
    COLOR_TEXTO_OSCURO = RGBColor(33, 33, 33)        # Texto principal
    COLOR_TEXTO_CLARO = RGBColor(255, 255, 255)      # Texto sobre fondos oscuros
    COLOR_FONDO_SUAVE = RGBColor(248, 249, 250)      # Fondo suave

    # === DIAPOSITIVA 1: PORTADA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fondo azul superior
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_AZUL_PRINCIPAL
    shape.line.fill.background()

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "🛣️ CORREDORES MINEROS Y TURÍSTICOS"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO_CLARO
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "Análisis Estadístico y Geoespacial para Identificar\nCorredores Complementarios a la Red Logística del MTC"
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_TEXTO_CLARO
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.3

    # Caja destacada
    key_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(3), Inches(8), Inches(0.4)
    )
    key_box.fill.solid()
    key_box.fill.fore_color.rgb = COLOR_NARANJA_ENFASIS
    key_box.line.fill.background()

    text_frame = key_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "⭐ Fortaleciendo la Competitividad Territorial con Inversión Pública y Privada"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO_CLARO
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Información del proyecto
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(2.8))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True

    info_items = [
        ("🏛️", "Programa: Buena Gobernanza Territorial – GIZ Perú"),
        ("👨‍💼", "Consultor: Jonatan Silvester Figueroa Gil"),
        ("📅", "Periodo: 28 octubre 2025 – 16 marzo 2026"),
        ("💰", "Valor: S/ 24,500")
    ]

    for i, (emoji, text) in enumerate(info_items):
        if i > 0:
            p = info_frame.add_paragraph()
        else:
            p = info_frame.paragraphs[0]
        p.text = f"{emoji} {text}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)

    # === DIAPOSITIVA 2: AGENDA ===
    slide = crear_slide_titulo(prs, "📋 Agenda de la Reunión Kick-off", COLOR_AZUL_PRINCIPAL)

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    agenda_items = [
        "1️⃣  Propósito de la Reunión",
        "2️⃣  El Contexto: ¿Por qué Corredores Complementarios?",
        "3️⃣  OBJETIVO PRINCIPAL de la Consultoría",
        "4️⃣  Objetivos Específicos",
        "5️⃣  Actividades Detalladas (10 actividades)",
        "6️⃣  Enfoque Metodológico Integrado",
        "7️⃣  Estructura de Trabajo y Entregables",
        "8️⃣  Cronograma Detallado",
        "9️⃣  Roles y Gestión del Proyecto",
        "🔟  Próximos Pasos"
    ]

    for i, item in enumerate(agenda_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.level = 0

    # === DIAPOSITIVA 3: PROPÓSITO ===
    slide = crear_slide_titulo(prs, "1️⃣ Propósito de la Reunión", COLOR_VERDE_EXITO)

    propositos = [
        ("🎯", "Alinear objetivos, alcance y resultados",
         "Asegurar entendimiento común entre GIZ y contrapartes técnicas (MEF, MTC, MINEM, MINCETUR)"),
        ("✅", "Validar la metodología integrada",
         "Confirmar el enfoque estadístico, geoespacial y econométrico propuesto"),
        ("📋", "Aprobar plan de trabajo y cronograma",
         "Revisar entregables P1 y P2 con sus fechas de entrega"),
        ("🤝", "Establecer gobernanza del proyecto",
         "Definir canales de comunicación y calendario de validaciones")
    ]

    y_pos = 2.2
    for emoji, titulo, descripcion in propositos:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.95)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_FONDO_SUAVE
        shape.line.color.rgb = COLOR_VERDE_EXITO
        shape.line.width = Pt(2.5)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.12)
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.2)

        p = text_frame.paragraphs[0]
        p.text = f"{emoji}  {titulo}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_VERDE_EXITO
        p.alignment = PP_ALIGN.LEFT

        p = text_frame.add_paragraph()
        p.text = descripcion
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2

        y_pos += 1.1

    # === DIAPOSITIVA 4: CONTEXTO/PROBLEMA ===
    slide = crear_slide_titulo(prs, "2️⃣ El Contexto: ¿Por qué Corredores Complementarios?", COLOR_ROJO_PROBLEMA)

    # Situación actual
    situacion_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.2)
    )
    situacion_box.fill.solid()
    situacion_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    situacion_box.line.color.rgb = COLOR_ROJO_PROBLEMA
    situacion_box.line.width = Pt(3)

    text_frame = situacion_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_left = Inches(0.3)

    p = text_frame.paragraphs[0]
    p.text = "⚠️ SITUACIÓN ACTUAL"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROJO_PROBLEMA
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.add_paragraph()
    p.text = "• La inversión pública se prioriza por brechas básicas sin enfoque territorial\n• El MTC definió corredores logísticos generales\n• Falta identificar corredores especializados para minería y turismo"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.line_spacing = 1.3
    p.space_before = Pt(8)

    # Oportunidad
    oportunidad_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(3.6), Inches(8.4), Inches(1.1)
    )
    oportunidad_box.fill.solid()
    oportunidad_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    oportunidad_box.line.color.rgb = COLOR_VERDE_EXITO
    oportunidad_box.line.width = Pt(3)

    text_frame = oportunidad_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_left = Inches(0.3)

    p = text_frame.paragraphs[0]
    p.text = "💡 OPORTUNIDAD"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE_EXITO
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.add_paragraph()
    p.text = "• Identificar corredores mineros y turísticos complementarios\n• Atraer inversión privada con enfoque territorial estratégico\n• Fortalecer competitividad regional"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.line_spacing = 1.3
    p.space_before = Pt(8)

    # Diagrama visual
    crear_diagrama_contexto(slide, COLOR_AZUL_PRINCIPAL, COLOR_VERDE_EXITO, COLOR_NARANJA_ENFASIS)

    # === DIAPOSITIVA 5: OBJETIVO PRINCIPAL ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fondo completo
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_AZUL_PRINCIPAL
    shape.line.fill.background()

    # Título de sección
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "3️⃣ OBJETIVO PRINCIPAL DE LA CONSULTORÍA"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXTO_CLARO
    p.alignment = PP_ALIGN.CENTER

    # Caja GRANDE y destacada para el objetivo principal
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.3), Inches(8.4), Inches(3.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 250, 230)
    shape.line.color.rgb = COLOR_NARANJA_ENFASIS
    shape.line.width = Pt(5)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.3)
    text_frame.margin_left = Inches(0.4)
    text_frame.margin_right = Inches(0.4)

    p = text_frame.paragraphs[0]
    p.text = "🎯 OBJETIVO GENERAL"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_NARANJA_ENFASIS
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

    p = text_frame.add_paragraph()
    p.text = "Elaborar un análisis estadístico y geoespacial que permita identificar corredores mineros y turísticos complementarios a los corredores logísticos existentes del MTC, con alto potencial de inversión pública y privada, mediante rutinas reproducibles de procesamiento de datos que fortalezcan la competitividad territorial."
    p.font.size = Pt(19)
    p.font.color.rgb = COLOR_TEXTO_OSCURO
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4

    # Elementos clave en la parte inferior
    y_pos = 5.8
    elementos_clave = [
        ("🛣️", "Corredores\nComplementarios", COLOR_AZUL_PRINCIPAL),
        ("⛏️", "Minería", COLOR_VERDE_EXITO),
        ("🏔️", "Turismo", COLOR_MORADO_ACCION),
        ("💼", "Inversión\nPública/Privada", COLOR_NARANJA_ENFASIS)
    ]

    x_pos = 1.2
    for emoji, texto, color in elementos_clave:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos), Inches(1.8), Inches(0.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        text_frame = box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = f"{emoji}\n{texto}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXTO_CLARO
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1

        x_pos += 2

    # === DIAPOSITIVA 6: OBJETIVOS ESPECÍFICOS ===
    slide = crear_slide_titulo(prs, "4️⃣ Objetivos Específicos", COLOR_VERDE_EXITO)

    objetivos_esp = [
        ("📊", "a)", "Recopilar, depurar y estandarizar bases de datos estadísticas, económicas, sociales y geoespaciales relevantes para el desarrollo minero y turístico"),
        ("📐", "b)", "Definir criterios técnicos y metodológicos para la identificación y priorización de corredores mineros y turísticos que complementen la red logística nacional"),
        ("🔍", "c)", "Procesar y analizar información estadística y econométrica que permita evaluar el potencial de economías de aglomeración y atracción de capital privado"),
        ("💻", "d)", "Desarrollar rutinas y algoritmos reproducibles (Python, R u otro software) para garantizar la actualización futura de la información"),
        ("🗺️", "e)", "Generar productos cartográficos y visualizaciones que orienten a los Gobiernos Regionales en la toma de decisiones")
    ]

    y_pos = 2.1
    for emoji, letra, objetivo in objetivos_esp:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.82)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_FONDO_SUAVE
        shape.line.color.rgb = COLOR_VERDE_EXITO
        shape.line.width = Pt(2.5)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.2)

        p = text_frame.paragraphs[0]
        p.text = f"{emoji}  {letra} {objetivo}"
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2

        y_pos += 0.95

    # === DIAPOSITIVA 7: ACTIVIDADES DETALLADAS (Parte 1) - FORMATO CORREGIDO ===
    slide = crear_slide_titulo(prs, "5️⃣ Actividades Detalladas (1/2)", COLOR_MORADO_ACCION)

    actividades_1 = [
        ("a)", "Sostener una reunión inicial (kick off) con el equipo del Proyecto Buena Gobernanza Territorial"),
        ("b)", "Revisar la propuesta conceptual del enfoque para la priorización territorial de la inversión pública"),
        ("c)", "Revisión, depuración y sistematización de bases de datos estadísticas y geoespaciales (INEI, MEF, MINEM, MINCETUR, OSINERGMIN, MINAM, GORE, MTC, etc.)"),
        ("d)", "Diseño metodológico para la identificación de corredores mineros y turísticos, asegurando complementariedad con los corredores logísticos nacionales del MTC"),
        ("e)", "Procesamiento estadístico y econométrico de información: limpieza, normalización, análisis de consistencia y estimaciones")
    ]

    y_pos = 2.2
    for letra, actividad in actividades_1:
        # UNA SOLA CAJA con todo el contenido para evitar separación
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.75)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(250, 245, 255)
        box.line.color.rgb = COLOR_MORADO_ACCION
        box.line.width = Pt(2)

        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.2)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        # FORMATO CORREGIDO: letra y actividad en la misma línea
        p.text = f"{letra} {actividad}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15

        y_pos += 0.95

    # === DIAPOSITIVA 8: ACTIVIDADES DETALLADAS (Parte 2) - FORMATO CORREGIDO ===
    slide = crear_slide_titulo(prs, "5️⃣ Actividades Detalladas (2/2)", COLOR_MORADO_ACCION)

    actividades_2 = [
        ("f)", "Procesamiento geoespacial: elaboración de mapas temáticos y capas SIG con identificación de corredores e inversiones habilitantes"),
        ("g)", "Identificar tipología de inversiones habilitantes para los corredores mineros y turísticos, con incidencia en regiones con mayor canon minero y principales activos turísticos"),
        ("h)", "Desarrollo de rutinas/algoritmos reproducibles para actualización de datos y mapas"),
        ("i)", "Presentación y validación de resultados preliminares con actores clave"),
        ("j)", "Elaboración del informe final con recomendaciones estratégicas")
    ]

    y_pos = 2.2
    for letra, actividad in actividades_2:
        # UNA SOLA CAJA con todo el contenido para evitar separación
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.75)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(250, 245, 255)
        box.line.color.rgb = COLOR_MORADO_ACCION
        box.line.width = Pt(2)

        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.2)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        # FORMATO CORREGIDO: letra y actividad en la misma línea
        p.text = f"{letra} {actividad}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15

        y_pos += 0.95

    # === DIAPOSITIVA 9: ENFOQUE METODOLÓGICO ===
    slide = crear_slide_titulo(prs, "6️⃣ Enfoque Metodológico Integrado", COLOR_NARANJA_ENFASIS)

    # Crear tabla
    rows, cols = 5, 3
    left = Inches(0.5)
    top = Inches(2.2)
    width = Inches(9)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Configurar anchos de columna
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.8)
    table.columns[2].width = Inches(3)

    # Encabezados
    headers = ["📋 Componente", "🔧 Metodología / Herramienta", "🎯 Resultado Esperado"]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NARANJA_ENFASIS
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_CLARO
        paragraph.alignment = PP_ALIGN.CENTER

    # Datos
    data = [
        ["📈 Estadístico-econométrico",
         "Normalización, correlaciones, K-means/DBSCAN, índices compuestos (IPC)",
         "Identificación de tipologías de corredores"],
        ["🗺️ Geoespacial",
         "Buffers 5 km, accesibilidad, LISA/Moran's I, mapas SIG",
         "Mapas temáticos y capas priorizadas"],
        ["💾 Datos y reproducibilidad",
         "Pipeline ETL, control QA, versión Git, metadata",
         "Repositorio digital reproducible"],
        ["⚖️ Transversalización",
         "Género e interculturalidad – criterios transformadores",
         "Indicadores desagregados por territorio"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(250, 250, 250) if i % 2 == 0 else RGBColor(255, 255, 255)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = COLOR_TEXTO_OSCURO
            paragraph.alignment = PP_ALIGN.LEFT
            cell.text_frame.word_wrap = True

    # === DIAPOSITIVA 10: ENTREGABLES ===
    slide = crear_slide_titulo(prs, "7️⃣ Estructura de Trabajo y Entregables", COLOR_AZUL_PRINCIPAL)

    # Timeline visual
    crear_timeline_entregables(slide, COLOR_AZUL_PRINCIPAL, COLOR_VERDE_EXITO, COLOR_NARANJA_ENFASIS)

    # === DIAPOSITIVA 11: CRONOGRAMA DETALLADO ===
    slide = crear_slide_titulo(prs, "8️⃣ Cronograma Detallado de Actividades", COLOR_VERDE_EXITO)

    # Crear tabla de cronograma
    rows, cols = 5, 3
    left = Inches(0.5)
    top = Inches(2.2)
    width = Inches(9)
    height = Inches(4.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Configurar anchos
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(2.5)

    # Encabezados
    headers = ["📅 Fase", "📋 Actividades Principales", "🎯 Producto/Resultado"]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_VERDE_EXITO
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_CLARO
        paragraph.alignment = PP_ALIGN.CENTER

    # Datos del cronograma
    cronograma_data = [
        ["F1 – Inicio\n28 oct-17 nov",
         "• Kick-off con contrapartes\n• Inventario de bases\n• Diseño metadata y Git\n• Criterios de priorización\n• Rutinas ETL y QA",
         "📄 Producto 1\nPlan de trabajo"],
        ["F2 – Procesamiento\nnov-dic 2025",
         "• Depuración de bases\n• Indicadores productividad\n• Modelos econométricos\n• Clusterización\n• Cálculo IPC y mapas SIG",
         "📊 Resultados intermedios y mapas"],
        ["F3 – Validación\nene 2026",
         "• Validación interna GIZ\n• Taller técnico sectorial\n• Ajustes metodológicos\n• Definición corredores priorizados",
         "📄 Producto 2\nInforme Final preliminar"],
        ["F4 – Cierre\nfeb-mar 2026",
         "• Consolidación base final\n• Repositorio digital\n• Manual de uso\n• Lecciones aprendidas",
         "✅ Informe de cierre y transferencia"]
    ]

    for i, row_data in enumerate(cronograma_data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = cell_text
            cell.fill.solid()
            colores_fila = [
                RGBColor(230, 240, 255),
                RGBColor(230, 250, 240),
                RGBColor(255, 245, 230),
                RGBColor(240, 230, 255)
            ]
            cell.fill.fore_color.rgb = colores_fila[i-1]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = COLOR_TEXTO_OSCURO
            paragraph.alignment = PP_ALIGN.LEFT
            cell.text_frame.word_wrap = True
            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_left = Inches(0.1)

    # === DIAPOSITIVA 12: ROLES Y COORDINACIÓN ===
    slide = crear_slide_titulo(prs, "9️⃣ Roles y Coordinación del Proyecto", COLOR_NARANJA_ENFASIS)

    roles = [
        ("🏢", "GIZ – Proyecto BGT", "Supervisión técnica y validación de productos (Jessica Ocsas)"),
        ("👨‍💼", "Consultor (J. Figueroa)", "Diseño metodológico, procesamiento, análisis y coordinación técnica"),
        ("🏛️", "MEF / MTC / MINEM / MINCETUR", "Contrapartes sectoriales – validación técnica y acceso a bases"),
        ("🌎", "GORE y actores regionales", "Validación territorial de corredores priorizados")
    ]

    y_pos = 2.3
    for emoji, titulo, descripcion in roles:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(y_pos), Inches(8), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 250, 240)
        shape.line.color.rgb = COLOR_NARANJA_ENFASIS
        shape.line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_left = Inches(0.2)

        p = text_frame.paragraphs[0]
        p.text = f"{emoji}  {titulo}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_NARANJA_ENFASIS

        p = text_frame.add_paragraph()
        p.text = descripcion
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.line_spacing = 1.2

        y_pos += 1.05

    # === DIAPOSITIVA 13: PRÓXIMOS PASOS ===
    slide = crear_slide_titulo(prs, "🔟 Próximos Pasos", COLOR_AZUL_PRINCIPAL)

    pasos = [
        ("📅 17 nov 2025", "Validar el Plan de Trabajo (P1)"),
        ("💾", "Configurar repositorio reproducible y diccionario de datos"),
        ("👥 Enero 2026", "Planificar taller técnico de validación con contrapartes"),
        ("📄 30 ene 2026", "Entregar Informe Final (P2)"),
        ("✅ 16 mar 2026", "Cierre contractual y liquidación")
    ]

    y_pos = 2.3
    for i, (emoji, paso) in enumerate(pasos):
        # Número de paso
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2), Inches(y_pos), Inches(0.4), Inches(0.4)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_AZUL_PRINCIPAL
        num_box.line.fill.background()

        num_frame = num_box.text_frame
        num_frame.text = str(i+1)
        num_frame.paragraphs[0].font.size = Pt(16)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = COLOR_TEXTO_CLARO
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Contenido del paso
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.8), Inches(y_pos-0.05), Inches(7), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 240, 255)
        shape.line.color.rgb = COLOR_AZUL_PRINCIPAL
        shape.line.width = Pt(1.5)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = f"{emoji}  {paso}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO_OSCURO

        y_pos += 0.75

    # === DIAPOSITIVA 14: CIERRE ===
    slide = crear_slide_cierre(prs, COLOR_AZUL_PRINCIPAL, COLOR_VERDE_EXITO, COLOR_NARANJA_ENFASIS)

    # Guardar presentación
    output_path = '/home/user/giz/Presentacion_Corredores_NARRATIVA_MEJORADA.pptx'
    prs.save(output_path)
    print(f"✅ Presentación creada exitosamente: {output_path}")
    return output_path

def crear_slide_titulo(prs, titulo, color):
    """Crear slide con título estilizado"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Barra superior de color
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    return slide

def crear_diagrama_contexto(slide, color_azul, color_verde, color_naranja):
    """Crear diagrama de contexto simplificado"""

    # MTC existente
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5), Inches(2.5), Inches(1)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(230, 245, 255)
    box1.line.color.rgb = color_azul
    box1.line.width = Pt(2)

    text_frame = box1.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = "🛣️ Corredores\nLogísticos MTC"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color_azul
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2

    # Símbolo +
    plus_box = slide.shapes.add_textbox(Inches(4.2), Inches(5.2), Inches(0.6), Inches(0.6))
    plus_frame = plus_box.text_frame
    p = plus_frame.paragraphs[0]
    p.text = "+"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color_naranja
    p.alignment = PP_ALIGN.CENTER

    # Nuevos corredores
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(5), Inches(2.5), Inches(1)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(255, 245, 230)
    box2.line.color.rgb = color_naranja
    box2.line.width = Pt(2)

    text_frame = box2.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = "⛏️🏔️ Corredores\nMineros y Turísticos"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color_naranja
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2

    # Flecha hacia resultado
    arrow_box = slide.shapes.add_textbox(Inches(4.5), Inches(6.2), Inches(1), Inches(0.5))
    arrow_frame = arrow_box.text_frame
    p = arrow_frame.paragraphs[0]
    p.text = "="
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color_verde
    p.alignment = PP_ALIGN.CENTER

    # Resultado
    result_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(6.5), Inches(6), Inches(0.7)
    )
    result_box.fill.solid()
    result_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    result_box.line.color.rgb = color_verde
    result_box.line.width = Pt(3)

    text_frame = result_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = "✅ Red Integrada de Corredores con Competitividad Territorial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color_verde
    p.alignment = PP_ALIGN.CENTER

def crear_timeline_entregables(slide, color_azul, color_verde, color_naranja):
    """Crear timeline visual de entregables"""

    # Línea de tiempo base
    slide.shapes.add_connector(
        1, Inches(1.5), Inches(4), Inches(8.5), Inches(4)
    ).line.color.rgb = RGBColor(150, 150, 150)

    # Entregables
    entregables = [
        (2, "📄 P1", "Plan de Trabajo\ny Avance Inicial", "17 nov 2025", color_verde),
        (5, "📊 P2", "Informe Final\nCompleto", "30 ene 2026", color_azul),
        (8, "✅", "Cierre\nContractual", "16 mar 2026", color_naranja)
    ]

    for x, emoji, titulo, fecha, color in entregables:
        # Círculo en timeline
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x-0.2), Inches(3.8), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Cuadro de información arriba
        info_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x-0.8), Inches(2.2), Inches(1.6), Inches(1.3)
        )
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        info_box.line.color.rgb = color
        info_box.line.width = Pt(2)

        text_frame = info_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        p = text_frame.add_paragraph()
        p.text = titulo
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1

        # Fecha abajo
        fecha_box = slide.shapes.add_textbox(
            Inches(x-0.8), Inches(4.5), Inches(1.6), Inches(0.4)
        )
        fecha_frame = fecha_box.text_frame
        p = fecha_frame.paragraphs[0]
        p.text = fecha
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    # Información de valor
    valor_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(5.5), Inches(4), Inches(0.8)
    )
    valor_box.fill.solid()
    valor_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    valor_box.line.color.rgb = color_verde
    valor_box.line.width = Pt(2)

    text_frame = valor_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "💰 Valor Total del Contrato: S/ 24,500"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color_verde
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def crear_slide_cierre(prs, color_azul, color_verde, color_naranja):
    """Crear slide de cierre con mensaje final"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fondo azul
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color_azul
    shape.line.fill.background()

    # Caja central con mensaje
    mensaje_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.8), Inches(8), Inches(4)
    )
    mensaje_box.fill.solid()
    mensaje_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    mensaje_box.line.fill.background()

    text_frame = mensaje_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.3)
    text_frame.margin_left = Inches(0.4)
    text_frame.margin_right = Inches(0.4)

    p = text_frame.paragraphs[0]
    p.text = "🎯 Impacto Esperado"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color_azul
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

    p = text_frame.add_paragraph()
    mensaje = """Este proyecto permitirá identificar corredores mineros y turísticos complementarios a la red logística del MTC, proporcionando evidencia estadística y geoespacial de alta calidad para orientar inversiones públicas y privadas hacia territorios con mayor potencial productivo, garantizando transparencia, reproducibilidad y enfoque territorial."""
    p.text = mensaje
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.5

    # Elementos clave
    p = text_frame.add_paragraph()
    p.text = ""
    p.space_after = Pt(15)

    elementos = "🛣️ Corredores Complementarios  |  ⛏️ Minería  |  🏔️ Turismo  |  💼 Inversión"
    p = text_frame.add_paragraph()
    p.text = elementos
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color_verde
    p.alignment = PP_ALIGN.CENTER

    # Mensaje de agradecimiento
    gracias_box = slide.shapes.add_textbox(Inches(2), Inches(6.3), Inches(6), Inches(0.8))
    gracias_frame = gracias_box.text_frame
    p = gracias_frame.paragraphs[0]
    p.text = "¡Gracias por su atención! 🤝"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide

if __name__ == "__main__":
    print("🚀 Iniciando creación de presentación con narrativa mejorada...")
    print("📝 Mejoras implementadas:")
    print("   ✓ Narrativa clara: Problema → Solución → Método → Resultados")
    print("   ✓ Formato de actividades corregido (sin saltos de línea)")
    print("   ✓ Colores con mejor contraste y legibilidad")
    print("   ✓ Jerarquía visual mejorada")
    crear_presentacion()
    print("✨ ¡Proceso completado!")
