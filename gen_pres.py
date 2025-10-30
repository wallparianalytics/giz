#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Colores
GIZ_BLUE = RGBColor(0, 51, 141)
GIZ_ORANGE = RGBColor(230, 126, 34)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(127, 127, 127)

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GIZ_BLUE
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'Identificación de Corredores\nMineros y Turísticos'
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = 'Análisis Estadístico y Geoespacial para la\nPriorización de Inversiones Públicas'
    p.font.size = Pt(22)
    p.font.color.rgb = GIZ_ORANGE
    p.alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = 'Proyecto: Buena Gobernanza Territorial - GIZ Perú'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p = tb.text_frame.add_paragraph()
    p.text = 'Consultor: Jonatan Silvester Figueroa Gil'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p = tb.text_frame.add_paragraph()
    p.text = 'Período: Octubre 2025 - Marzo 2026'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GIZ_BLUE

    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GIZ_ORANGE
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(8), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(15)
        p.level = 0

# SLIDE 1: Portada
add_title_slide(prs)

# SLIDE 2: Contexto
add_content_slide(prs, 'Contexto y Problemática', [
    '• Situación Actual: La inversión pública en Perú se prioriza por brechas básicas, sin enfoque territorial estratégico',
    '• Desafío: Falta información procesada y georreferenciada para identificar intervenciones estratégicas',
    '• Oportunidad: La Política Nacional de Inversión Pública impulsa el desarrollo de corredores económicos',
    '• Complementariedad: El MTC definió corredores logísticos. Se requiere identificar corredores mineros y turísticos'
])

# SLIDE 3: Propósito
add_content_slide(prs, 'Propósito de la Reunión', [
    '1. Alinear objetivos, alcance y resultados del estudio con GIZ y contrapartes técnicas (MEF, MTC, MINEM, MINCETUR)',
    '2. Validar la metodología integrada (estadística + geoespacial + econométrica)',
    '3. Aprobar el plan de trabajo, cronograma y productos (P1 y P2)',
    '4. Establecer la gobernanza del proyecto, canales de comunicación y calendario de validaciones'
])

# SLIDE 4: Objetivo General
slide = prs.slides.add_slide(prs.slide_layouts[6])
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = 'Objetivo General'
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = GIZ_BLUE

line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = GIZ_ORANGE
line.line.fill.background()

box = slide.shapes.add_shape(1, Inches(1), Inches(2), Inches(8), Inches(3))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(240, 248, 255)
box.line.color.rgb = GIZ_BLUE
box.line.width = Pt(3)

tf = box.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = 'Elaborar un análisis estadístico y geoespacial que permita identificar corredores mineros y turísticos complementarios a los del MTC, mediante rutinas reproducibles de procesamiento de datos y criterios de competitividad territorial.'
p.font.size = Pt(22)
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.4

tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'IMPACTO: Orientar inversiones públicas y privadas hacia territorios con mayor potencial de desarrollo, atracción de capital y generación de oportunidades.'
p.font.size = Pt(15)
p.font.italic = True
p.font.color.rgb = GIZ_ORANGE
p.alignment = PP_ALIGN.CENTER

# SLIDE 5: Objetivos Específicos
add_content_slide(prs, 'Objetivos Específicos', [
    '1. Recopilar, depurar y estandarizar bases estadísticas, económicas y geoespaciales',
    '2. Definir criterios metodológicos para priorización de corredores',
    '3. Evaluar potencial de aglomeración y atracción de capital privado',
    '4. Desarrollar algoritmos reproducibles en Python/R',
    '5. Elaborar mapas SIG y visualizaciones interactivas'
])

# SLIDE 6: Metodología
add_content_slide(prs, 'Enfoque Metodológico Integrado', [
    '📊 Estadístico-econométrico: K-means, DBSCAN, IPC → Tipologías de corredores',
    '🗺️ Geoespacial: Buffers 5km, LISA, Moran\'s I → Mapas temáticos priorizados',
    '💾 Datos y reproducibilidad: Pipeline ETL, Git, QA → Repositorio digital',
    '🤝 Transversalización: Género e interculturalidad → Indicadores desagregados'
])

# SLIDE 7: Cronograma
add_content_slide(prs, 'Cronograma de Actividades', [
    'F1 - Inicio y Planificación (Oct-Nov 2025): Plan de trabajo y avance inicial',
    'F2 - Procesamiento y Análisis (Nov-Dic 2025): Resultados intermedios y mapas',
    'F3 - Validación Técnica (Ene 2026): Informe Final preliminar',
    'F4 - Cierre y Transferencia (Feb-Mar 2026): Transferencia y documentación'
])

# SLIDE 8: Productos
slide = prs.slides.add_slide(prs.slide_layouts[6])
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = 'Productos Entregables'
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = GIZ_BLUE

line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = GIZ_ORANGE
line.line.fill.background()

# P1
box1 = slide.shapes.add_shape(1, Inches(0.8), Inches(1.8), Inches(4), Inches(4.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(240, 248, 255)
box1.line.color.rgb = GIZ_BLUE
box1.line.width = Pt(3)

tf = box1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'PRODUCTO 1'
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GIZ_BLUE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = 'Plan de Trabajo y Avance Inicial'
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(10)
p.space_after = Pt(15)

p = tf.add_paragraph()
p.text = '📅 17 noviembre 2025'
p.font.size = Pt(14)
p.font.color.rgb = GIZ_ORANGE
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(15)

for item in ['Cronograma detallado', 'Rutinas base (Python/R)', 'Mapas preliminares', 'Criterios de priorización']:
    p = tf.add_paragraph()
    p.text = f'✓ {item}'
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(5)

# P2
box2 = slide.shapes.add_shape(1, Inches(5.2), Inches(1.8), Inches(4), Inches(4.5))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(255, 250, 240)
box2.line.color.rgb = GIZ_ORANGE
box2.line.width = Pt(3)

tf = box2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'PRODUCTO 2'
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GIZ_ORANGE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = 'Informe Final Completo'
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(10)
p.space_after = Pt(15)

p = tf.add_paragraph()
p.text = '📅 30 enero 2026'
p.font.size = Pt(14)
p.font.color.rgb = GIZ_BLUE
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(15)

for item in ['Base integrada con metadatos', 'Mapas SIG definitivos', 'Repositorio reproducible', 'Documentación técnica']:
    p = tf.add_paragraph()
    p.text = f'✓ {item}'
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(5)

# SLIDE 9: Roles
add_content_slide(prs, 'Roles y Coordinación', [
    '• GIZ - Proyecto BGT: Supervisión técnica y validación (Jessica Ocsas)',
    '• Consultor (J. Figueroa): Diseño metodológico, procesamiento y análisis de datos',
    '• MEF / MTC / MINEM / MINCETUR: Contrapartes sectoriales - validación técnica',
    '• GORE y actores regionales: Validación territorial de corredores priorizados'
])

# SLIDE 10: Calidad y Riesgos
add_content_slide(prs, 'Gestión de Calidad y Riesgos', [
    '✓ Plan de Gestión de Datos (PGD) | Control QA | Conservación 10 años | PI cedida a GIZ',
    '⚠️ Retrasos en datos → Fuentes alternativas (INEI, Open Data MEF)',
    '⚠️ Inconsistencias de información → QA automático + verificación regional',
    '⚠️ Cambios metodológicos → Comité técnico quincenal',
    '⚠️ Fuerza mayor → Aplicación de cláusulas CCG'
])

# SLIDE 11: Impacto
add_content_slide(prs, 'Impacto Esperado', [
    '💡 Evidencia Estadística: Datos de alta calidad para orientar inversiones hacia territorios con potencial productivo',
    '🗺️ Mapas y Visualizaciones: Herramientas geoespaciales para priorización territorial y toma de decisiones',
    '🔄 Reproducibilidad: Rutinas automatizadas que permiten actualización continua de información',
    '🤝 Articulación Territorial: Complementariedad con corredores logísticos MTC para desarrollo integrado',
    '📊 Competitividad Regional: Identificación de corredores con mayor potencial de aglomeración',
    '🎯 Enfoque Transversal: Indicadores desagregados con perspectiva de género e interculturalidad'
])

# SLIDE 12: Próximos Pasos
add_content_slide(prs, 'Próximos Pasos', [
    '1. NOV: Validar Plan de Trabajo (P1) - 17 noviembre 2025',
    '2. NOV-DIC: Configurar repositorio reproducible - Implementación continua',
    '3. ENE: Planificar taller técnico de validación - Enero 2026',
    '4. ENE: Entregar Informe Final (P2) - 30 enero 2026',
    '5. MAR: Cierre y liquidación del proyecto - 16 marzo 2026'
])

# SLIDE 13: Cierre
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = GIZ_BLUE
bg.line.fill.background()

tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
tf = tb.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = '"Este proyecto permitirá a la GIZ y al MEF contar con evidencia estadística y geoespacial de alta calidad para orientar las inversiones públicas hacia territorios con mayor potencial productivo y turístico, garantizando transparencia, reproducibilidad y enfoque territorial."'
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.5

tb = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
p = tb.text_frame.paragraphs[0]
p.text = '¡Gracias por su atención!'
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = GIZ_ORANGE
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.7))
p = tb.text_frame.paragraphs[0]
p.text = 'Consultor: Jonatan Silvester Figueroa Gil | jonatan.figueroa.gil@gmail.com'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

prs.save('/home/user/giz/Presentacion_Corredores_Mineros_Turisticos.pptx')
print('✅ Presentación creada exitosamente!')
print('📄 13 slides generadas')
print('📁 Archivo: Presentacion_Corredores_Mineros_Turisticos.pptx')
