#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para crear presentación mejorada sobre Corredores Mineros y Turísticos
con gráficos, diagramas y emojis profesionales
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from io import BytesIO
import os

def crear_presentacion():
    """Crear presentación profesional completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Colores corporativos GIZ
    COLOR_AZUL_GIZ = RGBColor(0, 74, 151)
    COLOR_VERDE_GIZ = RGBColor(0, 135, 104)
    COLOR_NARANJA = RGBColor(255, 127, 39)
    COLOR_GRIS_OSCURO = RGBColor(51, 51, 51)
    COLOR_GRIS_CLARO = RGBColor(240, 240, 240)

    # === DIAPOSITIVA 1: PORTADA ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Fondo azul superior
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_AZUL_GIZ
    shape.line.fill.background()

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "📊 Procesamiento y Análisis de Inversiones Públicas"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Corredores Mineros y Turísticos Complementarios"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Información del proyecto
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(3))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True

    info_text = """🏛️ Programa: Buena Gobernanza Territorial – GIZ Perú

👨‍💼 Consultor: Jonatan Silvester Figueroa Gil

📅 Periodo: 28 octubre 2025 – 16 marzo 2026

💰 Valor: S/ 24,500
"""
    p = info_frame.paragraphs[0]
    p.text = info_text
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_GRIS_OSCURO
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(10)

    # === DIAPOSITIVA 2: AGENDA ===
    slide = crear_slide_titulo(prs, "📋 Agenda de la Reunión", COLOR_AZUL_GIZ)

    # Contenido
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    agenda_items = [
        ("1️⃣", "Propósito de la Reunión"),
        ("2️⃣", "Contexto del Proyecto"),
        ("3️⃣", "Objetivos General y Específicos"),
        ("4️⃣", "Enfoque Metodológico Integrado"),
        ("5️⃣", "Estructura de Trabajo y Entregables"),
        ("6️⃣", "Cronograma Detallado"),
        ("7️⃣", "Roles y Coordinación"),
        ("8️⃣", "Gestión de Calidad y Datos"),
        ("9️⃣", "Riesgos y Mitigación"),
        ("🔟", "Próximos Pasos")
    ]

    for emoji, item in agenda_items:
        p = text_frame.add_paragraph()
        p.text = f"{emoji}  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_GRIS_OSCURO
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

    # === DIAPOSITIVA 3: PROPÓSITO ===
    slide = crear_slide_titulo(prs, "1️⃣ Propósito de la Reunión", COLOR_VERDE_GIZ)

    propositos = [
        ("🎯", "Alinear objetivos, alcance y resultados",
         "Coordinación con GIZ y contrapartes técnicas (MEF, MTC, MINEM, MINCETUR)"),
        ("✅", "Validar la metodología integrada",
         "Enfoque estadístico + geoespacial + econométrico"),
        ("📑", "Aprobar plan de trabajo y cronograma",
         "Definición de productos P1 y P2 con fechas claras"),
        ("🤝", "Establecer gobernanza del proyecto",
         "Canales de comunicación y calendario de validaciones")
    ]

    y_pos = 2.2
    for emoji, titulo, detalle in propositos:
        agregar_cuadro_contenido(slide, emoji, titulo, detalle, y_pos, COLOR_VERDE_GIZ)
        y_pos += 1.1

    # === DIAPOSITIVA 4: CONTEXTO ===
    slide = crear_slide_titulo(prs, "2️⃣ Contexto del Proyecto", COLOR_NARANJA)

    # Diagrama de contexto
    crear_diagrama_contexto(slide)

    # === DIAPOSITIVA 5: OBJETIVO GENERAL ===
    slide = crear_slide_titulo(prs, "3️⃣ Objetivo General", COLOR_AZUL_GIZ)

    # Caja destacada para objetivo
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2.5), Inches(8), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 240, 255)
    shape.line.color.rgb = COLOR_AZUL_GIZ
    shape.line.width = Pt(3)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.2)
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)

    p = text_frame.paragraphs[0]
    p.text = "🎯 OBJETIVO GENERAL"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_AZUL_GIZ
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(15)

    p = text_frame.add_paragraph()
    p.text = "Elaborar un análisis estadístico y geoespacial que permita identificar corredores mineros y turísticos complementarios a los del MTC, mediante rutinas reproducibles de procesamiento de datos y criterios de competitividad territorial."
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_GRIS_OSCURO
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.3

    # === DIAPOSITIVA 6: OBJETIVOS ESPECÍFICOS ===
    slide = crear_slide_titulo(prs, "4️⃣ Objetivos Específicos", COLOR_VERDE_GIZ)

    objetivos_esp = [
        ("📊", "Recopilar, depurar y estandarizar bases estadísticas, económicas y geoespaciales"),
        ("📐", "Definir criterios metodológicos para priorización de corredores"),
        ("🔍", "Evaluar potencial de aglomeración y atracción de capital privado"),
        ("💻", "Desarrollar algoritmos reproducibles (Python/R)"),
        ("🗺️", "Elaborar mapas SIG y visualizaciones interactivas")
    ]

    y_pos = 2.2
    for emoji, objetivo in objetivos_esp:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_pos), Inches(7), Inches(0.7)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_GRIS_CLARO
        shape.line.color.rgb = COLOR_VERDE_GIZ
        shape.line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_left = Inches(0.2)

        p = text_frame.paragraphs[0]
        p.text = f"{emoji}  {objetivo}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_GRIS_OSCURO
        p.alignment = PP_ALIGN.LEFT

        y_pos += 0.85

    # === DIAPOSITIVA 7: ENFOQUE METODOLÓGICO (TABLA) ===
    slide = crear_slide_titulo(prs, "5️⃣ Enfoque Metodológico Integrado", COLOR_NARANJA)

    # Crear tabla
    rows, cols = 5, 3
    left = Inches(0.5)
    top = Inches(2.2)
    width = Inches(9)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Configurar anchos de columna
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.8)
    table.columns[2].width = Inches(3)

    # Encabezados
    headers = ["📋 Componente", "🔧 Metodología / Herramienta", "🎯 Resultado Esperado"]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NARANJA
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Datos
    data = [
        ["📈 Estadístico-econométrico",
         "Normalización, correlaciones, K-means/DBSCAN, índices compuestos (IPC)",
         "Identificación de tipologías de corredores"],
        ["🗺️ Geoespacial",
         "Buffers 5 km, accesibilidad, LISA/Moran's I, mapas SIG",
         "Mapas temáticos y capas priorizadas"],
        ["💾 Datos y reproducibilidad",
         "Pipeline ETL, control QA, versión Git, metadata",
         "Repositorio digital reproducible"],
        ["⚖️ Transversalización",
         "Género e interculturalidad – criterios transformadores",
         "Indicadores desagregados por territorio"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(250, 250, 250) if i % 2 == 0 else RGBColor(255, 255, 255)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = COLOR_GRIS_OSCURO
            paragraph.alignment = PP_ALIGN.LEFT
            cell.text_frame.word_wrap = True

    # === DIAPOSITIVA 8: ENTREGABLES ===
    slide = crear_slide_titulo(prs, "6️⃣ Estructura de Trabajo y Entregables", COLOR_AZUL_GIZ)

    # Timeline visual
    crear_timeline_entregables(slide)

    # === DIAPOSITIVA 9: CRONOGRAMA DETALLADO ===
    slide = crear_slide_titulo(prs, "7️⃣ Cronograma Detallado de Actividades", COLOR_VERDE_GIZ)

    # Crear tabla de cronograma
    rows, cols = 5, 3
    left = Inches(0.5)
    top = Inches(2.2)
    width = Inches(9)
    height = Inches(4.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Configurar anchos
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(2.5)

    # Encabezados
    headers = ["📅 Fase", "📋 Actividades Principales", "🎯 Producto/Resultado"]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_VERDE_GIZ
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Datos del cronograma
    cronograma_data = [
        ["F1 – Inicio\n28 oct-17 nov",
         "• Kick-off con contrapartes\n• Inventario de bases\n• Diseño metadata y Git\n• Criterios de priorización\n• Rutinas ETL y QA",
         "📄 Producto 1\nPlan de trabajo"],
        ["F2 – Procesamiento\nnov-dic 2025",
         "• Depuración de bases\n• Indicadores productividad\n• Modelos econométricos\n• Clusterización\n• Cálculo IPC y mapas SIG",
         "📊 Resultados intermedios y mapas"],
        ["F3 – Validación\nene 2026",
         "• Validación interna GIZ\n• Taller técnico sectorial\n• Ajustes metodológicos\n• Definición corredores priorizados",
         "📄 Producto 2\nInforme Final preliminar"],
        ["F4 – Cierre\nfeb-mar 2026",
         "• Consolidación base final\n• Repositorio digital\n• Manual de uso\n• Lecciones aprendidas",
         "✅ Informe de cierre y transferencia"]
    ]

    for i, row_data in enumerate(cronograma_data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = cell_text
            cell.fill.solid()
            colores_fila = [
                RGBColor(230, 240, 255),
                RGBColor(230, 250, 240),
                RGBColor(255, 245, 230),
                RGBColor(240, 230, 255)
            ]
            cell.fill.fore_color.rgb = colores_fila[i-1]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = COLOR_GRIS_OSCURO
            paragraph.alignment = PP_ALIGN.LEFT
            cell.text_frame.word_wrap = True
            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_left = Inches(0.1)

    # === DIAPOSITIVA 10: ROLES Y COORDINACIÓN ===
    slide = crear_slide_titulo(prs, "8️⃣ Roles y Coordinación", COLOR_NARANJA)

    roles = [
        ("🏢 GIZ – Proyecto BGT", "Supervisión técnica y validación de productos (Jessica Ocsas)"),
        ("👨‍💼 Consultor (J. Figueroa)", "Diseño metodológico, procesamiento, análisis y coordinación técnica"),
        ("🏛️ MEF / MTC / MINEM / MINCETUR", "Contrapartes sectoriales – validación técnica y acceso a bases"),
        ("🌎 GORE y actores regionales", "Validación territorial de corredores priorizados")
    ]

    y_pos = 2.3
    for emoji, descripcion in roles:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(y_pos), Inches(8), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 250, 240)
        shape.line.color.rgb = COLOR_NARANJA
        shape.line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_left = Inches(0.2)

        p = text_frame.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(18)
        p.font.bold = True

        p = text_frame.add_paragraph()
        p.text = descripcion
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_GRIS_OSCURO
        p.line_spacing = 1.2

        y_pos += 1.05

    # === DIAPOSITIVA 11: GESTIÓN DE CALIDAD ===
    slide = crear_slide_titulo(prs, "9️⃣ Gestión de Calidad y Datos", COLOR_AZUL_GIZ)

    calidad_items = [
        ("📋", "Plan de Gestión de Datos (PGD)",
         "Confidencialidad, backup, versionado y control de calidad"),
        ("✅", "Control de Calidad Metodológico",
         "Validación cruzada de indicadores y mapas"),
        ("📦", "Conservación Documental",
         "Resguardo por 10 años conforme a CCG locales GIZ"),
        ("©️", "Propiedad Intelectual",
         "Derechos de autor cedidos a GIZ")
    ]

    y_pos = 2.3
    for emoji, titulo, detalle in calidad_items:
        agregar_cuadro_contenido(slide, emoji, titulo, detalle, y_pos, COLOR_AZUL_GIZ)
        y_pos += 1.15

    # === DIAPOSITIVA 12: RIESGOS Y MITIGACIÓN ===
    slide = crear_slide_titulo(prs, "🔟 Riesgos y Mitigación", COLOR_VERDE_GIZ)

    # Crear tabla de riesgos
    rows, cols = 5, 2
    left = Inches(0.8)
    top = Inches(2.2)
    width = Inches(8.4)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Configurar anchos
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4.4)

    # Encabezados
    headers = ["⚠️ Riesgo", "🛡️ Estrategia de Mitigación"]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_VERDE_GIZ
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(13)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Datos de riesgos
    riesgos_data = [
        ["Retrasos en acceso a datos oficiales",
         "Plan de contingencia con fuentes alternativas (INEI, Open Data MEF)"],
        ["Inconsistencias de información territorial",
         "QA automático + verificación manual regional"],
        ["Cambios metodológicos por contrapartes",
         "Comité técnico quincenal para aprobación de ajustes"],
        ["Eventos de fuerza mayor (bioseguridad)",
         "Aplicación de cláusulas CCG y ajuste de plazos"]
    ]

    for i, row_data in enumerate(riesgos_data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.rows[i].cells[j]
            cell.text = cell_text
            cell.fill.solid()
            if j == 0:
                cell.fill.fore_color.rgb = RGBColor(255, 240, 240)
            else:
                cell.fill.fore_color.rgb = RGBColor(240, 255, 240)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = COLOR_GRIS_OSCURO
            paragraph.alignment = PP_ALIGN.LEFT
            cell.text_frame.word_wrap = True
            cell.text_frame.margin_top = Inches(0.1)
            cell.text_frame.margin_left = Inches(0.15)

    # === DIAPOSITIVA 13: PRÓXIMOS PASOS ===
    slide = crear_slide_titulo(prs, "1️⃣1️⃣ Próximos Pasos", COLOR_NARANJA)

    pasos = [
        ("📅 17 nov 2025", "Validar el Plan de Trabajo (P1)"),
        ("💾", "Configurar repositorio reproducible y diccionario de datos"),
        ("👥 Enero 2026", "Planificar taller técnico de validación con contrapartes"),
        ("📄 30 ene 2026", "Entregar Informe Final (P2)"),
        ("✅ 16 mar 2026", "Cierre contractual y liquidación")
    ]

    y_pos = 2.3
    for i, (emoji, paso) in enumerate(pasos):
        # Número de paso
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.2), Inches(y_pos), Inches(0.4), Inches(0.4)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_NARANJA
        num_box.line.fill.background()

        num_frame = num_box.text_frame
        num_frame.text = str(i+1)
        num_frame.paragraphs[0].font.size = Pt(16)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Contenido del paso
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.8), Inches(y_pos-0.05), Inches(7), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 250, 240)
        shape.line.color.rgb = COLOR_NARANJA
        shape.line.width = Pt(1.5)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = f"{emoji}  {paso}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_GRIS_OSCURO

        y_pos += 0.75

    # === DIAPOSITIVA 14: CIERRE ===
    slide = crear_slide_cierre(prs, COLOR_AZUL_GIZ)

    # Guardar presentación
    output_path = '/home/user/giz/Presentacion_Corredores_Mineros_Turisticos_MEJORADA.pptx'
    prs.save(output_path)
    print(f"✅ Presentación creada exitosamente: {output_path}")
    return output_path

def crear_slide_titulo(prs, titulo, color):
    """Crear slide con título estilizado"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Barra superior de color
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    return slide

def agregar_cuadro_contenido(slide, emoji, titulo, detalle, y_pos, color):
    """Agregar cuadro de contenido con emoji"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(y_pos), Inches(7), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_left = Inches(0.2)

    p = text_frame.paragraphs[0]
    p.text = f"{emoji}  {titulo}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color

    p = text_frame.add_paragraph()
    p.text = detalle
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.line_spacing = 1.1

def crear_diagrama_contexto(slide):
    """Crear diagrama visual del contexto"""
    COLOR_AZUL = RGBColor(0, 74, 151)
    COLOR_VERDE = RGBColor(0, 135, 104)
    COLOR_NARANJA = RGBColor(255, 127, 39)

    # Cuadros de contexto
    contextos = [
        (2, 2.5, "⚠️ SITUACIÓN ACTUAL",
         "Inversión pública priorizada\npor brechas básicas sin\nenfoque territorial",
         RGBColor(255, 230, 230)),
        (5.5, 2.5, "🛣️ AVANCE MTC",
         "Definición de corredores\nlogísticos nacionales",
         RGBColor(230, 245, 255)),
        (2, 4.5, "🎯 NECESIDAD",
         "Identificar corredores\nmineros y turísticos\ncomplementarios",
         RGBColor(255, 245, 230)),
        (5.5, 4.5, "💼 OBJETIVO",
         "Atraer inversión privada\ny articular desarrollo\nregional",
         RGBColor(230, 255, 230))
    ]

    for x, y, titulo, texto, color_fondo in contextos:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(3), Inches(1.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_fondo
        shape.line.color.rgb = COLOR_AZUL
        shape.line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_top = Inches(0.1)

        p = text_frame.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_AZUL
        p.alignment = PP_ALIGN.CENTER

        p = text_frame.add_paragraph()
        p.text = texto
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2

    # Flechas de conexión
    # Flecha 1: Situación -> Necesidad
    agregar_flecha(slide, 3.5, 3.9, 3.5, 4.5, COLOR_NARANJA)

    # Flecha 2: Avance -> Objetivo
    agregar_flecha(slide, 7, 3.9, 7, 4.5, COLOR_VERDE)

def agregar_flecha(slide, x1, y1, x2, y2, color):
    """Agregar flecha entre dos puntos"""
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(3)

def crear_timeline_entregables(slide):
    """Crear timeline visual de entregables"""
    COLOR_AZUL = RGBColor(0, 74, 151)
    COLOR_VERDE = RGBColor(0, 135, 104)
    COLOR_NARANJA = RGBColor(255, 127, 39)

    # Línea de tiempo base
    slide.shapes.add_connector(
        1, Inches(1.5), Inches(4), Inches(8.5), Inches(4)
    ).line.color.rgb = RGBColor(150, 150, 150)

    # Entregables
    entregables = [
        (2, "📄 P1", "Plan de Trabajo\ny Avance Inicial", "17 nov 2025", COLOR_VERDE),
        (5, "📊 P2", "Informe Final\nCompleto", "30 ene 2026", COLOR_AZUL),
        (8, "✅", "Cierre\nContractual", "16 mar 2026", COLOR_NARANJA)
    ]

    for x, emoji, titulo, fecha, color in entregables:
        # Círculo en timeline
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x-0.2), Inches(3.8), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Cuadro de información arriba
        info_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x-0.8), Inches(2.2), Inches(1.6), Inches(1.3)
        )
        info_box.fill.solid()
        info_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
        info_box.line.color.rgb = color
        info_box.line.width = Pt(2)

        text_frame = info_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        p = text_frame.add_paragraph()
        p.text = titulo
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1

        # Fecha abajo
        fecha_box = slide.shapes.add_textbox(
            Inches(x-0.8), Inches(4.5), Inches(1.6), Inches(0.4)
        )
        fecha_frame = fecha_box.text_frame
        p = fecha_frame.paragraphs[0]
        p.text = fecha
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    # Información de valor
    valor_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(5.5), Inches(4), Inches(0.8)
    )
    valor_box.fill.solid()
    valor_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
    valor_box.line.color.rgb = COLOR_VERDE
    valor_box.line.width = Pt(2)

    text_frame = valor_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "💰 Valor Total del Contrato: S/ 24,500"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_VERDE
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def crear_slide_cierre(prs, color):
    """Crear slide de cierre con mensaje final"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fondo degradado simulado con rectángulos
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Caja central con mensaje
    mensaje_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2), Inches(8), Inches(3.5)
    )
    mensaje_box.fill.solid()
    mensaje_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    mensaje_box.line.fill.background()

    text_frame = mensaje_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.3)
    text_frame.margin_left = Inches(0.4)
    text_frame.margin_right = Inches(0.4)

    p = text_frame.paragraphs[0]
    p.text = "1️⃣2️⃣ Cierre del Kick-off"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

    p = text_frame.add_paragraph()
    mensaje = """🎯 Este proyecto permitirá a la GIZ y al MEF contar con evidencia estadística y geoespacial de alta calidad para orientar las inversiones públicas hacia territorios con mayor potencial productivo y turístico, garantizando transparencia, reproducibilidad y enfoque territorial."""
    p.text = mensaje
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4

    # Mensaje de agradecimiento
    gracias_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.8))
    gracias_frame = gracias_box.text_frame
    p = gracias_frame.paragraphs[0]
    p.text = "¡Gracias por su atención! 🤝"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide

if __name__ == "__main__":
    print("🚀 Iniciando creación de presentación mejorada...")
    crear_presentacion()
    print("✨ ¡Proceso completado!")
